"""
PowerPoint Converter for Jupyter Notebooks
Converts Jupyter notebooks to PowerPoint presentations using KULeuven template
"""

import base64
import io
import nbformat as nbf
from glob import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER, MSO_SHAPE_TYPE
from PIL import UnidentifiedImageError
from IPython.lib.latextools import latex_to_png
import copy
from pptx.shapes.autoshape import Shape
from bs4 import BeautifulSoup
import re

# Collect a list of all notebooks in the content folder
notebooks = glob("ToegepasteAnalogeElektronica/*.ipynb", recursive=True)
notebooks+= glob("AnalogeElektronica2/*.ipynb", recursive=True)
notebooks+= glob("AnalogDesignTechniques/*.ipynb", recursive=True)
#notebooks+= glob("MicroEnNanoTechnologie/*.ipynb", recursive=True)

#slidetemplate="./.github/common/KULeuventemplate.pptx"
#slidetemplate="./.github/common/KULeuven_600TEMPLATE.pptx"
slidetemplate="./.github/common/PowerPoint_600jaar_template.pptx"

KUL_layout_title=0
KUL_layout_subtitle=1
KUL_layout_outline=9
KUL_layout_fig=4
KUL_layout_code=3
KUL_layout_text=6
KUL_layout_mdtext=6
KUL_layout_statement=11
KUL_layout_end=10

lines_per_chunk=11


prs = Presentation(slidetemplate)
body_left, body_top, body_width, body_height=Inches(1), Inches(0.62), prs.slide_width-Inches(1.5), Inches(6.2)

report_all_shapes_in_template=True
if report_all_shapes_in_template:
    for i1,lay in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(lay)
        print(i1, len(slide.shapes))
        for i2,shap in enumerate(slide.shapes):
            if shap.shape_type==14:
                print("---",i2,shap.shape_type,shap.placeholder_format.type)
            else:
                print("---",i2,shap.shape_type) 
        for shape in slide.placeholders:
            print('.  %d %s' % (shape.placeholder_format.idx, shape.name))

def clone_shape(shape, left=body_left, top=body_top, width=body_width, height=body_height,idcounter=0):
    """Add a duplicate of `shape` to the slide on which it appears."""
    sp = shape._sp
    spTree = sp.getparent()
    new_sp = copy.deepcopy(sp)
    spTree.append(new_sp)
    new_shape = Shape(new_sp, None)
    #new_shape.shape_id = shape.shape_id + 1345 +idcounter # Ensure a unique shape ID. AttributeError: property 'shape_id' of 'Shape' object has no setter
    new_shape.left = left
    new_shape.top = top
    new_shape.width = width
    new_shape.height = height
    return new_shape
    """  van een andere bron: mogelijk dit gebruiken als we meer fragments willen kopieren
        new_el = copy.deepcopy(shape.element)
        spTree.shapes._spTree.insert_element_before(new_el, "p:extLst") # in de bron was de eerste spTree de besteming-slide
        new_shape = spTree.shapes[-1]  # in de bron was de eerste spTree de besteming-slide
    """

def find_between( s, first, last ):
    try:
        start = s.index( first )
        end = s.index( last, start+len(first) )+len(last)
        return s[start:end]
    except ValueError:
        return ""

def process_latex_in_text(text):
    """Process text and convert LaTeX symbols between $ $ to Unicode."""
    
    def replace_latex_match(match):
        latex_content = match.group(1)
        return latex_to_unicode(latex_content)
    
    return re.sub(r'\$(.*?)\$', replace_latex_match, text)

def maketitle(cell,slide):
    if "KULeuvenSlides" in cell.get('metadata', {}):
        if "slide_title" in cell.metadata.get('KULeuvenSlides', {}):
            st = process_latex_in_text(cell.metadata.KULeuvenSlides["slide_title"])
            slide.shapes.title.text=st.replace("<BR>","\n")
            slide.shapes.title.text_frame.fit_text(font_family="Arial",max_size=32,bold=True, font_file=r".github/common/fonts/arialbd.ttf")         

def parse_markdown_bullets(md_text):
    lines = md_text.strip().split('\n')
    bullets = []
    for line in lines:
        stripped = line.lstrip()
        if stripped.startswith(('-', '*', '+')):
            level = (len(line) - len(stripped)) // 2  # 2 spaces per indent
            bullets.append((level+1, stripped[1:].strip()))
        else:
            bullets.append((0, line))
    return bullets

def latex_to_unicode(latex_text):
    """Convert common LaTeX symbols to Unicode equivalents."""
    # Dictionary mapping LaTeX commands to Unicode symbols
    latex_unicode_map = {
        # Greek letters (lowercase)
        r'\alpha': 'α', r'\beta': 'β', r'\gamma': 'γ', r'\delta': 'δ',
        r'\epsilon': 'ε', r'\varepsilon': 'ε', r'\zeta': 'ζ', r'\eta': 'η',
        r'\theta': 'θ', r'\vartheta': 'ϑ', r'\iota': 'ι', r'\kappa': 'κ',
        r'\lambda': 'λ', r'\mu': 'μ', r'\nu': 'ν', r'\xi': 'ξ',
        r'\pi': 'π', r'\varpi': 'ϖ', r'\rho': 'ρ', r'\varrho': 'ϱ',
        r'\sigma': 'σ', r'\varsigma': 'ς', r'\tau': 'τ', r'\upsilon': 'υ',
        r'\phi': 'φ', r'\varphi': 'φ', r'\chi': 'χ', r'\psi': 'ψ', r'\omega': 'ω',
        
        # Greek letters (uppercase)
        r'\Alpha': 'Α', r'\Beta': 'Β', r'\Gamma': 'Γ', r'\Delta': 'Δ',
        r'\Epsilon': 'Ε', r'\Zeta': 'Ζ', r'\Eta': 'Η', r'\Theta': 'Θ',
        r'\Iota': 'Ι', r'\Kappa': 'Κ', r'\Lambda': 'Λ', r'\Mu': 'Μ',
        r'\Nu': 'Ν', r'\Xi': 'Ξ', r'\Pi': 'Π', r'\Rho': 'Ρ',
        r'\Sigma': 'Σ', r'\Tau': 'Τ', r'\Upsilon': 'Υ', r'\Phi': 'Φ',
        r'\Chi': 'Χ', r'\Psi': 'Ψ', r'\Omega': 'Ω',
        
        # Mathematical operators
        r'\pm': '±', r'\mp': '∓', r'\times': '×', r'\div': '÷',
        r'\cdot': '·', r'\bullet': '•', r'\circ': '∘', r'\ast': '∗',
        r'\star': '⋆', r'\dagger': '†', r'\ddagger': '‡', r'\cap': '∩',
        r'\cup': '∪', r'\uplus': '⊎', r'\sqcap': '⊓', r'\sqcup': '⊔',
        r'\vee': '∨', r'\wedge': '∧', r'\oplus': '⊕', r'\ominus': '⊖',
        r'\otimes': '⊗', r'\oslash': '⊘', r'\odot': '⊙', r'\bigcirc': '◯',
        r'\diamond': '⋄', r'\bigtriangleup': '△', r'\bigtriangledown': '▽',
        r'\triangleleft': '◁', r'\triangleright': '▷', r'\lhd': '⊲', r'\rhd': '⊳',
        r'\unlhd': '⊴', r'\unrhd': '⊵', r'\amalg': '⨿',
        
        # Relations
        r'\leq': '≤', r'\geq': '≥', r'\equiv': '≡', r'\models': '⊨',
        r'\prec': '≺', r'\succ': '≻', r'\sim': '∼', r'\perp': '⊥',
        r'\preceq': '⪯', r'\succeq': '⪰', r'\simeq': '≃', r'\mid': '∣',
        r'\ll': '≪', r'\gg': '≫', r'\asymp': '≍', r'\parallel': '∥',
        r'\subset': '⊂', r'\supset': '⊃', r'\approx': '≈', r'\bowtie': '⋈',
        r'\subseteq': '⊆', r'\supseteq': '⊇', r'\cong': '≅', r'\sqsubset': '⊏',
        r'\sqsupset': '⊐', r'\neq': '≠', r'\smile': '⌣', r'\sqsubseteq': '⊑',
        r'\sqsupseteq': '⊒', r'\doteq': '≐', r'\frown': '⌢', r'\in': '∈',
        r'\ni': '∋', r'\propto': '∝', r'\vdash': '⊢', r'\dashv': '⊣',
        
        # Arrows
        r'\leftarrow': '←', r'\gets': '←', r'\rightarrow': '→', r'\to': '→',
        r'\leftrightarrow': '↔', r'\uparrow': '↑', r'\downarrow': '↓',
        r'\updownarrow': '↕', r'\Leftarrow': '⇐', r'\Rightarrow': '⇒',
        r'\Leftrightarrow': '⇔', r'\Uparrow': '⇑', r'\Downarrow': '⇓',
        r'\Updownarrow': '⇕', r'\mapsto': '↦', r'\longmapsto': '⟼',
        r'\hookleftarrow': '↩', r'\hookrightarrow': '↪', r'\leftharpoonup': '↼',
        r'\leftharpoondown': '↽', r'\rightharpoonup': '⇀', r'\rightharpoondown': '⇁',
        r'\rightleftharpoons': '⇌', r'\leadsto': '⇝',
        
        # Miscellaneous symbols
        r'\infty': '∞', r'\aleph': 'ℵ', r'\hbar': 'ℏ', r'\imath': 'ı',
        r'\jmath': 'ȷ', r'\ell': 'ℓ', r'\wp': '℘', r'\Re': 'ℜ', r'\Im': 'ℑ',
        r'\mho': '℧', r'\prime': '′', r'\emptyset': '∅', r'\nabla': '∇',
        r'\surd': '√', r'\partial': '∂', r'\top': '⊤', r'\bot': '⊥',
        r'\vdots': '⋮', r'\ddots': '⋱', r'\heartsuit': '♡', r'\diamondsuit': '♢',
        r'\clubsuit': '♣', r'\spadesuit': '♠', r'\neg': '¬', r'\flat': '♭',
        r'\natural': '♮', r'\sharp': '♯',
        
        # Mathematical functions
        r'\sin': 'sin', r'\cos': 'cos', r'\tan': 'tan', r'\cot': 'cot',
        r'\sec': 'sec', r'\csc': 'csc', r'\arcsin': 'arcsin', r'\arccos': 'arccos',
        r'\arctan': 'arctan', r'\sinh': 'sinh', r'\cosh': 'cosh', r'\tanh': 'tanh',
        r'\coth': 'coth', r'\exp': 'exp', r'\log': 'log', r'\ln': 'ln',
        r'\det': 'det', r'\gcd': 'gcd', r'\lcm': 'lcm', r'\lim': 'lim',
        r'\liminf': 'lim inf', r'\limsup': 'lim sup', r'\max': 'max', r'\min': 'min',
        r'\sup': 'sup', r'\inf': 'inf', r'\arg': 'arg', r'\ker': 'ker',
        r'\dim': 'dim', r'\hom': 'hom', r'\deg': 'deg',
        
        # Large operators
        r'\sum': '∑', r'\prod': '∏', r'\coprod': '∐', r'\int': '∫',
        r'\oint': '∮', r'\iint': '∬', r'\iiint': '∭', r'\iiiint': '⨌',
        r'\idotsint': '∫⋯∫', r'\bigcup': '⋃', r'\bigcap': '⋂',
        r'\biguplus': '⨄', r'\bigsqcup': '⨆', r'\bigvee': '⋁',
        r'\bigwedge': '⋀', r'\bigodot': '⨀', r'\bigotimes': '⨂',
        r'\bigoplus': '⨁', r'\bigcirc': '◯',
        
        # Delimiters
        r'\langle': '⟨', r'\rangle': '⟩', r'\lceil': '⌈', r'\rceil': '⌉',
        r'\lfloor': '⌊', r'\rfloor': '⌋', r'\{': '{', r'\}': '}',
        
        # Subscripts and superscripts (common ones)
        r'_0': '₀', r'_1': '₁', r'_2': '₂', r'_3': '₃', r'_4': '₄',
        r'_5': '₅', r'_6': '₆', r'_7': '₇', r'_8': '₈', r'_9': '₉',
        r'^0': '⁰', r'^1': '¹', r'^2': '²', r'^3': '³', r'^4': '⁴',
        r'^5': '⁵', r'^6': '⁶', r'^7': '⁷', r'^8': '⁸', r'^9': '⁹',
        r'^+': '⁺', r'^-': '⁻', r'^=': '⁼', r'^(': '⁽', r'^)': '⁾',
        
        # Common fractions
        r'\frac{1}{2}': '½', r'\frac{1}{3}': '⅓', r'\frac{2}{3}': '⅔',
        r'\frac{1}{4}': '¼', r'\frac{3}{4}': '¾', r'\frac{1}{5}': '⅕',
        r'\frac{2}{5}': '⅖', r'\frac{3}{5}': '⅗', r'\frac{4}{5}': '⅘',
        r'\frac{1}{6}': '⅙', r'\frac{5}{6}': '⅚', r'\frac{1}{7}': '⅐',
        r'\frac{1}{8}': '⅛', r'\frac{3}{8}': '⅜', r'\frac{5}{8}': '⅝',
        r'\frac{7}{8}': '⅞', r'\frac{1}{9}': '⅑', r'\frac{1}{10}': '⅒',
    }
    
    result = latex_text
    # Sort by length (longest first) to avoid partial replacements
    for latex_cmd in sorted(latex_unicode_map.keys(), key=len, reverse=True):
        result = result.replace(latex_cmd, latex_unicode_map[latex_cmd])
    
    return result

def add_parsed_bullet(paragraph, text):
    
    # Combine all patterns into one regex with named groups
    combined_pattern = (
        r"(?P<sub><sub>(.*?)</sub>)|"
        r"(?P<sup><sup>(.*?)</sup>)|"
        r"(?P<bold>\*\*(.*?)\*\*)|"
        r"(?P<dollar>\$(.*?)\$)"
    )
    last_end = 0
    
    for match in re.finditer(combined_pattern, text):
        run = paragraph.add_run()
        start, end = match.span()
        # Add text before the tag (if any)
        if start > last_end:
            run.text = text[last_end:start]
            run = paragraph.add_run()
        # Determine which pattern matched
        if match.group("sub"):
            run.text = match.group(2)
            run.font.size = Pt(18)
            #run.font.subscript = True
            run.font._element.set('baseline', '-25000')
        elif match.group("sup"):
            run.text =  match.group(4)
            run.font.size = Pt(18)
            #run.font.superscript = True
            run.font._element.set('baseline', '45000')
        elif match.group("bold"):
            run.text = match.group(6)
            run.font.bold = True
        elif match.group("dollar"):
            run.text =  latex_to_unicode(match.group(8))
        last_end = end

    # Add any remaining text after the last tag
    if last_end < len(text):
        run = paragraph.add_run()
        run.text = text[last_end:]

for ipath in notebooks:
    print("file om te zetten: ",ipath)
    ntbk = nbf.read(ipath, nbf.NO_CONVERT)
    idcounter=0
    prs = Presentation(slidetemplate)
    slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_title])
    if "title" in ntbk.metadata.get('KULeuvenSlides', {}):
        title_shape = slide.shapes.title
        title_shape.text = ntbk.metadata.KULeuvenSlides["title"].replace("<BR>","\n")
        footer=ntbk.metadata.KULeuvenSlides["title"].replace("<BR>","\n")
        #title_shape.text_frame.fit_text()
        #title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        #title_shape.text_frame.paragraphs[0].font.bold = True
        #title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        #title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if "subtitle" in ntbk.metadata.get('KULeuvenSlides', {}):
        subtitle= slide.shapes.add_textbox(Inches(7),Inches(4), Inches(5),Inches(2.0))  #left, top, width, height
        subtitle.text = ntbk.metadata.KULeuvenSlides["subtitle"].replace("<BR>","\n")
        #subtitle.text_frame.fit_text()
        for par in subtitle.text_frame.paragraphs:
            par.font.size = Pt(24)
            par.font.color.rgb = (RGBColor(255, 255, 255))  # White text
        #subtitle.text_frame.paragraphs[0].alignment = (PP_ALIGN.CENTER)
    if "authors" in ntbk.metadata.get('KULeuvenSlides', {}):
        txBox = slide.shapes.add_textbox(Inches(7),Inches(6.1), Inches(2),Inches(0.5))  #left, top, width, height
        tf = txBox.text_frame
        tf.text = ntbk.metadata.KULeuvenSlides["authors"]
        #tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        #tf.paragraphs[0].alignment = PP_ALIGN.CENTER        
    if "date" in ntbk.metadata.get('KULeuvenSlides', {}):
        txBox = slide.shapes.add_textbox(Inches(7),Inches(6.3), Inches(2),Inches(0.5))  #left, top, width, height
        tf = txBox.text_frame
        tf.text = ntbk.metadata.KULeuvenSlides["date"]
        #tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        #tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for index,cell in enumerate(ntbk.cells):
        if "code" in cell.get('cell_type', {}) and not("remove_cell4pptx" in cell.metadata.get('tags', {})):
            if "slide_type" in cell.metadata.get('slideshow', {}):
                if  cell.metadata.slideshow.get("slide_type", ())=="slide":
                   for output_idx, output in enumerate(cell.get("outputs", ())):
                        if "image/png" in output.get("data", {}):                           
                            slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_fig])
                            maketitle(cell,slide)  
                            image_stream = io.BytesIO(base64.b64decode(output.data["image/png"]))
                            try:                            
                                pictp=slide.shapes.add_picture(image_stream, body_left, body_top, height=body_height) 
                                if pictp.width>prs.slide_width:
                                    factorsc=(body_height*prs.slide_width)//pictp.width
                                    pictp.width=prs.slide_width
                                    pictp.height= factorsc
                                    pictp.left=Inches(0)
                                else:
                                    pictp.left=(prs.slide_width-pictp.width)//2
                                running_height=body_top+pictp.height+Inches(0.1)
                            except UnidentifiedImageError:
                                print("  image/png  error for cell number "+str(index))
                                running_height=body_top
                        elif "image/jpeg" in output.get("data", {}):                           
                            slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_fig])
                            maketitle(cell,slide)  
                            image_stream = io.BytesIO(base64.b64decode(output.data["image/jpeg"]))
                            try:                            
                                pictp=slide.shapes.add_picture(image_stream, body_left, body_top, height=body_height) 
                                if pictp.width>prs.slide_width:
                                    factorsc=(body_height*prs.slide_width)//pictp.width
                                    pictp.width=prs.slide_width
                                    pictp.height= factorsc
                                    pictp.left=Inches(0)
                                else:
                                    pictp.left=(prs.slide_width-pictp.width)//2
                                running_height=body_top+pictp.height+Inches(0.1)
                            except UnidentifiedImageError:
                                print("  image/jpeg  error for cell number "+str(index))
                                running_height=body_top
                        elif "image/gif" in output.get("data", {}):                           
                            slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_fig])
                            maketitle(cell,slide)  
                            image_stream = io.BytesIO(base64.b64decode(output.data["image/gif"]))
                            try:                            
                                pictp=slide.shapes.add_picture(image_stream, body_left, body_top, height=body_height) 
                                if pictp.width>prs.slide_width:
                                    factorsc=(body_height*prs.slide_width)//pictp.width
                                    pictp.width=prs.slide_width
                                    pictp.height= factorsc
                                    pictp.left=Inches(0)
                                else:
                                    pictp.left=(prs.slide_width-pictp.width)//2
                                running_height=body_top+pictp.height+Inches(0.1)
                            except UnidentifiedImageError:
                                print("  image/gif  error for cell number "+str(index))
                                running_height=body_top
                        elif "".join(cell.source).startswith("display.Video("):
                            # Extract video path from Video() call
                            video_path = None
                            source_line = cell.source.strip()
                            print(source_line)
                            # Look for Video("path") pattern
                            import re
                            video_match = re.search(r'display.Video\s*\(\s*["\']([^"\']+)["\']', source_line)
                            if video_match:
                                video_path = ipath.split('/')[0] + '/' + video_match.group(1)[2:]
                                print(f"  Found Video call with path: {video_path}")
                                
                                slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_fig])
                                maketitle(cell,slide)  
                                try:                            
                                    film=slide.shapes.add_movie(video_path, body_left, body_top,
                                                                 height=body_height,width=body_width,poster_frame_image=video_path+'.jpeg', mime_type='video/mp4')
                                    if film.width>prs.slide_width:
                                        film.width=prs.slide_width
                                        film.left=Inches(0)
                                    else:
                                        film.left=(prs.slide_width-film.width)//2
                                    running_height=body_top+film.height+Inches(0.1)
                                except UnidentifiedImageError:
                                    print("  video/mp4  error for cell number "+str(index))
                                    running_height=body_top
                                running_height = body_top + Inches(4)
                            else:
                                print(f"  Could not extract path from Video call in cell {index}")
                                running_height = body_top
                        
                        # elif "text/html" in output.get("data", {}):  # er zitten veel header lines voor de html
                            # lines="".join(output.data["text/html"]).split('\n')
                            # for isl in range(0,len(lines),lines_per_chunk):
                                # slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_code])
                                # maketitle(cell,slide)
                                # for isp in range(isl,min(isl+lines_per_chunk,len(lines))):
                                    # soup = BeautifulSoup(lines[isp]), "html.parser") ##. dit moet niet per lijn maar op het geheel
                                    # p=slide.shapes[0].text_frame.paragraphs[isp-isl] 
                                    # for span in soup.find_all("span"):
                                        # style = span.get("style", "")
                                        # color_match = re.search(r"color:\s*(#[0-9a-fA-F]{6})", style)
                                        # color = RGBColor(255, 255, 255)  # default white
                                        # if color_match:
                                            # hex_color = color_match.group(1)
                                            # color = RGBColor(*(int(hex_color[i:i+2], 16) for i in (1, 3, 5)))
                                        # run = p.add_run()
                                        # run.text = span.text
                                        # run.font.color.rgb = color
                                        # run.font.name = 'Courier New'   
                                        # run.font.size = Pt(14)
                        elif "text/plain" in output.get("data", {}):
                            lines="".join(output.data["text/plain"]).split('\n')
                            for i in range(0,len(lines),lines_per_chunk):
                                slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_code])
                                maketitle(cell,slide) 
                                slide.shapes[0].text="\n".join(lines[i:i+lines_per_chunk])
                                for par in slide.shapes[0].text_frame.paragraphs:
                                    par.line_spacing = Pt(24)
                                    par.font.color.rgb = RGBColor(255, 255, 255)
                                slide.shapes[0].text_frame.fit_text(font_family="Courier",max_size=24, font_file=r".github/common/fonts/cour.ttf")
                                running_height=body_top+slide.shapes[0].height+Inches(0.3)
                        elif "text" in output:
                            lines="".join(output["text"]).split('\n')
                            for i in range(0,len(lines),lines_per_chunk):
                                slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_text])
                                maketitle(cell,slide) 
                                slide.shapes[0].text="\n".join(lines[i:i+lines_per_chunk])
                                for par in slide.shapes[0].text_frame.paragraphs:
                                    par.line_spacing = Pt(24)
                                    par.font.color.rgb = RGBColor(0,0,0)
                                slide.shapes[0].text_frame.fit_text(font_family="Courier",max_size=24, font_file=r".github/common/fonts/cour.ttf") 
                                running_height=body_top+slide.shapes[0].height+Inches(0.3)         
                        else:
                            print(f"  content_type error for cell number {index}")
                            print(f"    Available output data types: {list(output.get('data', {}).keys())}")
                            print(f"    Available cell.outputs keys: {list(cell.outputs[output_idx].keys()) if hasattr(cell, 'outputs') and output_idx < len(cell.outputs) else 'N/A'}")
                            print(f"    All cell.outputs structure: {[list(out.keys()) for out in cell.outputs] if hasattr(cell, 'outputs') else 'No outputs'}")
                            
        if "markdown" in cell.get('cell_type', {}) and not("remove_cell4pptx" in cell.metadata.get('tags', {})):
            if "slide_type" in cell.metadata.get('slideshow', {}):
                md_text="".join(cell.get('source', {}))
                bullets = parse_markdown_bullets(md_text)
                if cell.metadata.slideshow.slide_type == "slide":
                    if ("KULeuvenSlides" in cell.get('metadata', {}) and "slide_code" in cell.metadata.get('KULeuvenSlides', {}) and cell.metadata.KULeuvenSlides.slide_code == "title"):
                        slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_subtitle])
                        maketitle(cell,slide)
                        running_height=body_top
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_mdtext])
                        maketitle(cell,slide)
                        body_shape = slide.shapes[0]
                        tf = body_shape.text_frame
                        tf.clear()  # Remove any existing paragraphs
                        tf.text = ""  # Clear existing text
                        for idx, (level, text) in enumerate(bullets):
                            if idx == 0:
                                # Reuse the existing paragraph for the first bullet
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.level = level
                            add_parsed_bullet(p, text)
                        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        running_height = body_top+body_shape.height + Inches(0.2)
                    
                elif cell.metadata.slideshow.get("slide_type", ())=="notes":
                    if slide is not None:  # Notes require a previous slide to exist
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = md_text

                # Process slide content for slide and fragment types
                horizontal_shift=Inches(0)
                scaling_factor=1.0
                if slide is not None and cell.metadata.slideshow.get("slide_type", ())=="fragment": 
                    if "KULeuvenSlides" in cell.get('metadata', {}):
                        if "eq_vertical" in cell.metadata.get('KULeuvenSlides', {}):
                            if cell.metadata.KULeuvenSlides["eq_vertical"]:
                                running_height+=Inches(cell.metadata.KULeuvenSlides["eq_vertical"])
                        if "eq_horizontal" in cell.metadata.get('KULeuvenSlides', {}):
                            if cell.metadata.KULeuvenSlides["eq_horizontal"]:
                                horizontal_shift=Inches(cell.metadata.KULeuvenSlides["eq_horizontal"])
                        if "eq_scale" in cell.metadata.get('KULeuvenSlides', {}):
                            if cell.metadata.KULeuvenSlides["eq_scale"]:
                                scaling_factor=cell.metadata.KULeuvenSlides["eq_scale"]
                                
                    latexpng=find_between(md_text  , "$$", "$$" )
                    latexpng2=find_between( md_text , r"\begin{equation}", r"\end{equation}" )
                    if len(latexpng2)>0:
                        latexpng=latexpng2
                    if len(latexpng)>0:
                        pictp = None
                        try:
                            pictp=slide.shapes.add_picture(io.BytesIO(latex_to_png(latexpng,backend="dvipng",scale=2*scaling_factor)), Inches(1),running_height) 
                        except UnidentifiedImageError:
                            print("   latex error for cell number "+str(index))
                        except Exception as e:
                            print(f"   latex error for cell number {index}: {e}")    
                        if pictp is not None:
                            if pictp.width>prs.slide_width:
                                factorsc=(body_height*prs.slide_width)//pictp.width
                                pictp.width=prs.slide_width
                                pictp.height= factorsc
                                pictp.left=horizontal_shift
                            else:
                                pictp.left=horizontal_shift+(prs.slide_width-pictp.width)//2
                            running_height+=pictp.height+Inches(0.3)
                    else:
                        body_shape = clone_shape(slide.shapes[0],  top=running_height, idcounter=idcounter)
                        idcounter+=1
                        tf = body_shape.text_frame
                        #tf.clear()  # Remove any existing paragraphs
                        #tf.text = ""  # Clear existing text
                        for idx, (level, text) in enumerate(bullets):
                            if idx == 0:
                                # Reuse the existing paragraph for the first bullet
                                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                            else:
                                p = tf.add_paragraph()
                            p.level = level
                            add_parsed_bullet(p, text)
                        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        running_height += body_shape.height + Inches(0.2)
                                   

# Loop through all slides except the first
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue

        # Loop through placeholders to find footer and slide number
        footer_found = False
        slide_number_found = False
        
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.idx == 10:  # Footer placeholder
                    shape.text = footer
                    footer_found = True
                elif shape.placeholder_format.idx == 11:  # Slide number placeholder
                    shape.text = f"{i + 1}"
                    slide_number_found = True
            except:
                # Skip if placeholder doesn't support text or has other issues
                continue
        
        # If standard placeholders weren't found, try alternative approaches
        if not footer_found or not slide_number_found:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Check if this might be a footer or slide number area based on position
                    # Footer typically at bottom, slide number typically at bottom right
                    slide_height = prs.slide_height
                    slide_width = prs.slide_width
                    
                    # If shape is in bottom 10% of slide
                    if shape.top > slide_height * 0.9:
                        if not footer_found and shape.left < slide_width * 0.5:  # Left side = footer
                            try:
                                shape.text = footer
                                footer_found = True
                            except:
                                pass
                        elif not slide_number_found and shape.left > slide_width * 0.7:  # Right side = slide number
                            try:
                                shape.text = f"{i + 1}"
                                slide_number_found = True
                            except:
                                pass

    slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_statement])
    slide = prs.slides.add_slide(prs.slide_layouts[KUL_layout_end])
    prs.save(ipath[:-6]+".pptx")